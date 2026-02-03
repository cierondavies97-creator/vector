"""
Deterministic indexing pipeline with:
- resume / skip unchanged files
- cooperative cancellation
- semantic chunking
- tier-1 + tier-2 + tier-3 tag enrichment
- detailed progress events (documents, chunks, embeddings, FAISS)
"""

from __future__ import annotations

import json
import hashlib
import threading

import time
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Callable
import re
from pdfminer.high_level import extract_text
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from memory_engine import IndexStats
from semantic_chunker import semantic_chunk
from embedding_enrichment import enrich_for_embedding
from concepts import infer_semantic_tags


# ============================================================
# Progress / Events
# ============================================================

class Stage(str, Enum):
    DISCOVER = "Discovering files"
    FILTER = "Filtering files"
    PARSE = "Parsing files"
    CHUNK = "Chunking documents"
    STORE = "Storing chunks"
    INDEX = "Embedding & indexing"


@dataclass
class ProgressEvent:
    stage: Stage
    current: int
    total: int
    message: str
    file: str | None = None
    debug: bool = False


class ProgressSink:
    def __init__(self, emit: Callable[[ProgressEvent], None]):
        self.emit = emit

    def stage(self, stage: Stage, total: int, message: str):
        self.emit(ProgressEvent(stage, 0, total, message))

    def step(self, stage: Stage, current: int, total: int, file: str | None = None):
        self.emit(
            ProgressEvent(
                stage=stage,
                current=current,
                total=total,
                message="Processing",
                file=file,
            )
        )

    def log(self, stage: Stage, message: str):
        self.emit(
            ProgressEvent(
                stage=stage,
                current=0,
                total=0,
                message=message,
                debug=True,
            )
        )


# ============================================================
# Data Models
# ============================================================

@dataclass(frozen=True)
class ParsedDocument:
    path: str
    text: str


@dataclass(frozen=True)
class Chunk:
    id: str
    source: str
    text: str
    metadata: dict


# ============================================================
# Indexing Pipeline
# ============================================================

class IndexingPipeline:
    SUPPORTED_SUFFIXES = {
        ".txt", ".md", ".pdf", ".py", ".json", ".docx", ".xlsx", ".pptx"
    }

    def __init__(
        self,
        memory_engine,
        kb_path: Path,
        chunk_size: int = 3000,
        overlap: int = 200,
    ):
        self.memory_engine = memory_engine
        self.kb_path = kb_path
        self.chunk_size = chunk_size
        self.overlap = overlap

        self.ingest_path = kb_path / "ingested"
        self.chunk_path = self.ingest_path / "chunks"
        self.meta_path = self.ingest_path / "meta.json"

        self.chunk_path.mkdir(parents=True, exist_ok=True)

        self._cancel_event = threading.Event()
        self.meta = self._load_meta()

    # --------------------------------------------------------
    # Control
    # --------------------------------------------------------

    def cancel(self):
        self._cancel_event.set()

    def _cancelled(self) -> bool:
        return self._cancel_event.is_set()

    # --------------------------------------------------------
    # Public Entry
    # --------------------------------------------------------

    def run(self, root: Path, sink: ProgressSink):
        files = self._discover(root, sink)
        if self._cancelled(): return

        files = self._filter(files, sink)
        if self._cancelled(): return

        docs = self._parse(files, sink)
        if self._cancelled(): return

        chunks = self._chunk(docs, sink)
        if self._cancelled(): return

        self._store(chunks, sink)
        if self._cancelled(): return

        stats = self._index(chunks, sink)
        self._save_meta()
        return stats

    # --------------------------------------------------------
    # Metadata
    # --------------------------------------------------------

    def _load_meta(self) -> dict:
        if self.meta_path.exists():
            try:
                return json.loads(self.meta_path.read_text())
            except Exception:
                pass
        return {"files": {}}

    def _save_meta(self):
        self.meta_path.write_text(json.dumps(self.meta, indent=2))

    def _file_hash(self, path: Path) -> str:
        h = hashlib.sha256()
        with path.open("rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()

    # --------------------------------------------------------
    # Discover / Filter
    # --------------------------------------------------------

    def _discover(self, root: Path, sink: ProgressSink) -> list[Path]:
        sink.stage(Stage.DISCOVER, 0, f"Scanning {root}")
        files = [p for p in root.rglob("*") if p.is_file() and not p.is_symlink()]
        sink.log(Stage.DISCOVER, f"Found {len(files)} files")
        return files

    def _filter(self, files: list[Path], sink: ProgressSink) -> list[Path]:
        accepted = []
        for p in files:
            if p.name.startswith("~$"):
                continue
            if p.suffix.lower() not in self.SUPPORTED_SUFFIXES:
                continue
            accepted.append(p)
        sink.log(Stage.FILTER, f"Accepted {len(accepted)} files")
        return accepted

    # --------------------------------------------------------
    # Parse
    # --------------------------------------------------------

    def _parse(self, files: list[Path], sink: ProgressSink) -> list[ParsedDocument]:
        docs: list[ParsedDocument] = []
        sink.stage(Stage.PARSE, len(files), "Parsing files")

        for i, path in enumerate(files, start=1):
            if self._cancelled():
                return docs

            sink.step(Stage.PARSE, i, len(files), path.name)

            digest = self._file_hash(path)
            mtime = int(path.stat().st_mtime)

            prev = self.meta["files"].get(path.as_posix())
            if prev and prev["hash"] == digest and prev["mtime"] == mtime:
                continue

            text = self._extract_text(path)
            if text.strip():
                docs.append(ParsedDocument(path=path.as_posix(), text=text))
                self.meta["files"][path.as_posix()] = {
                    "hash": digest,
                    "mtime": mtime,
                }

        sink.log(Stage.PARSE, f"Parsed {len(docs)} documents")
        return docs

    # --------------------------------------------------------
    # Chunking
    # --------------------------------------------------------

    def _chunk(self, docs: list[ParsedDocument], sink: ProgressSink) -> list[Chunk]:
        chunks: list[Chunk] = []
        sink.stage(Stage.CHUNK, len(docs), "Chunking documents")

        for i, doc in enumerate(docs, start=1):
            if self._cancelled():
                return chunks

            sink.step(Stage.CHUNK, i, len(docs), Path(doc.path).name)

            structured = semantic_chunk(
                doc.text,
                max_chars=self.chunk_size,
            )

            for j, sc in enumerate(structured):
                cid = f"{doc.path}:{j}"
                chunks.append(Chunk(cid, doc.path, sc["text"], sc.get("metadata", {})))

        sink.log(Stage.CHUNK, f"Created {len(chunks)} chunks")
        return chunks

    # --------------------------------------------------------
    # Store
    # --------------------------------------------------------

    def _store(self, chunks: list[Chunk], sink: ProgressSink):
        sink.stage(Stage.STORE, len(chunks), "Writing chunks")

        for i, c in enumerate(chunks, start=1):
            if self._cancelled():
                return

            safe = c.id.replace(":", "__").replace("/", "__").replace("\\", "__")
            (self.chunk_path / f"{safe}.txt").write_text(c.text, encoding="utf-8")
            metadata = self._build_chunk_metadata(c)
            (self.chunk_path / f"{safe}.meta.json").write_text(
                json.dumps(metadata, indent=2),
                encoding="utf-8",
            )

            if i % 250 == 0 or i == len(chunks):
                sink.step(Stage.STORE, i, len(chunks))

        sink.log(Stage.STORE, f"Stored {len(chunks)} chunks")

    # --------------------------------------------------------
    # Index (Embedding + FAISS)
    # --------------------------------------------------------

    def _index(self, chunks: list[Chunk], sink: ProgressSink) -> IndexStats:
        total = len(chunks)
        sink.stage(Stage.INDEX, total, "Embedding & indexing")

        mapping = {c.id: c.text for c in chunks}
        if self.memory_engine.config.embedding_enrichment_enabled:
            mapping = {
                cid: enrich_for_embedding(
                    text=text,
                    kind=self.memory_engine.config.embedding_enrichment_kind,
                    source=cid.split(":", 1)[0],
                )
                for cid, text in mapping.items()
            }

        # ---- embedding progress ----
        def progress_cb(cur: int, total: int, msg: str):
            sink.emit(
                ProgressEvent(
                    stage=Stage.INDEX,
                    current=cur,
                    total=total,
                    message=msg,
                )
            )

        t0 = time.perf_counter()
        stats = self.memory_engine.build_index(
            mapping,
            progress_cb=progress_cb,
        )
        t1 = time.perf_counter()

        faiss_seconds = max(t1 - t0 - stats.embedding_seconds, 0.0)

        sink.log(
            Stage.INDEX,
            f"Embedding done | chunks={stats.chunk_count} "
            f"embed={stats.embedding_seconds:.2f}s "
            f"faiss={faiss_seconds:.2f}s",
        )
        return IndexStats(
            embedding_seconds=stats.embedding_seconds,
            faiss_seconds=faiss_seconds,
            chunk_count=stats.chunk_count,
        )

    # --------------------------------------------------------
    # Text Extraction
    # --------------------------------------------------------

    def _extract_text(self, path: Path) -> str:
        suf = path.suffix.lower()

        if suf in {".txt", ".md"}:
            return path.read_text("utf-8", errors="ignore")

        if suf == ".pdf":
            return extract_text(path)

        if suf == ".docx":
            return "\n".join(p.text for p in Document(path).paragraphs)

        if suf == ".pptx":
            pres = Presentation(path)
            return "\n".join(
                shape.text
                for slide in pres.slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )

        if suf == ".xlsx":
            wb = load_workbook(path, read_only=True, data_only=True)
            lines = []
            for sh in wb.worksheets:
                for row in sh.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c]
                    if vals:
                        lines.append("\t".join(vals))
            return "\n".join(lines)

        if suf == ".json":
            try:
                return json.dumps(json.loads(path.read_text()), indent=2)
            except Exception:
                return path.read_text()

        return ""

    # --------------------------------------------------------
    # Metadata
    # --------------------------------------------------------

    def _infer_source_type(self, path: Path) -> str:
        parts = {p.lower() for p in path.parts}
        if {"peer_reviewed", "peer-reviewed", "journal", "papers", "paper"} & parts:
            return "peer_reviewed"
        if {"opinion", "opinions", "commentary", "blog", "editorial"} & parts:
            return "opinionated"
        if {"factual", "facts", "fact"} & parts:
            return "factual"
        return "general"

    def _build_chunk_metadata(self, chunk: Chunk) -> dict:
        source_path = Path(chunk.source)
        source_type = self._infer_source_type(source_path)
        doc_type = source_path.suffix.lstrip(".").lower() or "unknown"

        file_meta = self.meta["files"].get(source_path.as_posix(), {})
        source_hash = file_meta.get("hash")
        source_mtime = file_meta.get("mtime")

        chunk_index = int(chunk.id.rsplit(":", 1)[-1])
        is_overview = chunk_index == 0

        tags = [
            "namespace:file",
            f"source_type:{source_type}",
            f"doc_type:{doc_type}",
        ]

        heading = chunk.metadata.get("heading")
        heading_level = chunk.metadata.get("heading_level")
        section_path = chunk.metadata.get("section_path") or []

        tags.extend(
            infer_semantic_tags(
                text=chunk.text,
                heading=heading,
                section_path=section_path,
            )
        )

        return {
            "source_path": source_path.as_posix(),
            "chunk_index": chunk_index,
            "heading": heading,
            "heading_level": heading_level,
            "section_path": section_path,
            "source_type": source_type,
            "doc_type": doc_type,
            "source_hash": source_hash,
            "source_mtime": source_mtime,
            "is_overview": is_overview,
            "tags": tags,
            "factors": {
                "namespace": "file",
                "source_type": source_type,
                "doc_type": doc_type,
                "source_path": source_path.as_posix(),
                "chunk_index": chunk_index,
                "heading": heading,
                "heading_level": heading_level,
                "section_path": section_path,
            },
        }
