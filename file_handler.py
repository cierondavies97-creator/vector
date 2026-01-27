"""File ingestion and indexing utilities."""

from __future__ import annotations

import json
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Dict, Iterable

from pdfminer.high_level import extract_text
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from config import AppConfig
from memory_engine import MemoryEngine


@dataclass(frozen=True)
class ParsedFile:
    path: str
    content: str


@dataclass(frozen=True)
class IngestStats:
    processed: int
    skipped: int
    candidate_files: int
    parse_seconds: float
    embedding_seconds: float
    faiss_seconds: float
    chunk_count: int


class FileHandler:
    SUPPORTED_SUFFIXES = {
        ".txt", ".md", ".pdf", ".py", ".json", ".docx", ".xlsx", ".pptx"
    }

    def __init__(self, config: AppConfig, memory_engine: MemoryEngine) -> None:
        self.config = config
        self.memory_engine = memory_engine
        self.kb_path = self.config.knowledge_base_path() / "ingested"
        self.kb_path.mkdir(parents=True, exist_ok=True)

    # =========================
    # PUBLIC ENTRY POINT
    # =========================

    def reindex_workspace(
        self,
        workspace_path: str,
        progress_cb: ProgressCallback | None = None,
        debug_cb: DebugCallback | None = None,
    ) -> IngestStats:

        def debug(msg: str):
            if debug_cb:
                debug_cb(msg)

        root = Path(workspace_path)

        # 1. DISCOVER
        candidates = self._discover_files(root, debug)

        # 2. FILTER
        candidates = self._filter_files(candidates, debug)

        # 3. PARSE
        documents = self._parse_files(candidates, progress_cb, debug)

        # 4. CHUNK
        chunks = self._chunk_documents(documents, debug)

        # 5. STORE
        self._store_chunks(chunks, debug)

        # 6. INDEX
        index_stats = self._build_index(chunks, debug)

        return index_stats

    # =========================
    # STAGE 1 — DISCOVER
    # =========================

    def _discover_files(self, root: Path, debug) -> list[Path]:
        debug(f"[DISCOVER] Scanning {root}")
        files = [
            p for p in root.rglob("*")
            if p.is_file() and not p.is_symlink()
        ]
        debug(f"[DISCOVER] Found {len(files)} files")
        return files

    # =========================
    # STAGE 2 — FILTER
    # =========================

    def _filter_files(self, files: list[Path], debug) -> list[Path]:
        filtered = []

        max_bytes = self.config.max_file_mb * 1024 * 1024

        for p in files:
            if p.name.startswith("~$"):
                continue
            if p.name in {".DS_Store", "Thumbs.db", "desktop.ini"}:
                continue
            if p.suffix.lower() not in self.SUPPORTED_SUFFIXES:
                continue
            try:
                if p.stat().st_size > max_bytes:
                    continue
            except OSError:
                continue

            filtered.append(p)

        debug(f"[FILTER] Accepted {len(filtered)} files")
        return filtered

    # =========================
    # STAGE 3 — PARSE
    # =========================

    def _parse_files(
        self,
        files: list[Path],
        progress_cb,
        debug,
    ) -> list[ParsedFile]:

        parsed: list[ParsedFile] = []
        total = len(files)

        for idx, path in enumerate(files, start=1):
            if progress_cb:
                progress_cb(idx, total, f"Parsing {path.name}")

            try:
                content = self._extract_text(path)
                if content.strip():
                    parsed.append(
                        ParsedFile(
                            path=path.as_posix(),
                            content=content,
                        )
                    )
            except Exception as exc:
                debug(f"[PARSE] Failed {path}: {exc}")

        debug(f"[PARSE] Parsed {len(parsed)} documents")
        return parsed

    # =========================
    # STAGE 4 — CHUNK
    # =========================

    def _chunk_documents(self, docs: list[ParsedFile], debug) -> list[ParsedFile]:
        # NOTE: chunking is deferred to MemoryEngine for now
        # This keeps compatibility with your current design
        debug(f"[CHUNK] Passing {len(docs)} docs to indexer")
        return docs

    # =========================
    # STAGE 5 — STORE
    # =========================

    def _store_chunks(self, docs: list[ParsedFile], debug):
        for doc in docs:
            safe_name = (
                doc.path.replace(":", "")
                .replace("\\", "__")
                .replace("/", "__")
            )
            target = self.kb_path / f"{safe_name}.txt"
            target.write_text(doc.content, encoding="utf-8")

        debug(f"[STORE] Stored {len(docs)} files")

    # =========================
    # STAGE 6 — INDEX
    # =========================

    def _build_index(self, docs: list[ParsedFile], debug) -> IngestStats:
        debug("[INDEX] Building vector index")

        mapping = {doc.path: doc.content for doc in docs}
        stats = self.memory_engine.build_index(mapping)

        debug(
            f"[INDEX] chunks={stats.chunk_count} "
            f"embed={stats.embedding_seconds:.2f}s "
            f"faiss={stats.faiss_seconds:.2f}s"
        )

        return IngestStats(
            processed=len(docs),
            skipped=0,
            candidate_files=len(docs),
            parse_seconds=0.0,
            embedding_seconds=stats.embedding_seconds,
            faiss_seconds=stats.faiss_seconds,
            chunk_count=stats.chunk_count,
        )

    # =========================
    # FILE-TYPE DISPATCH
    # =========================

    def _extract_text(self, path: Path) -> str:
        suffix = path.suffix.lower()

        if suffix in {".txt", ".md"}:
            return path.read_text("utf-8", errors="ignore")

        if suffix == ".json":
            raw = path.read_text("utf-8", errors="ignore")
            try:
                return json.dumps(json.loads(raw), indent=2)
            except json.JSONDecodeError:
                return raw

        if suffix == ".pdf":
            return extract_text(path)

        if suffix == ".docx":
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        if suffix == ".xlsx":
            wb = load_workbook(path, read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    row_vals = ["" if c is None else str(c) for c in row]
                    if any(v.strip() for v in row_vals):
                        lines.append("\t".join(row_vals))
            return "\n".join(lines)

        if suffix == ".pptx":
            pres = Presentation(path)
            lines = []
            for i, slide in enumerate(pres.slides, 1):
                lines.append(f"# Slide {i}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
            return "\n".join(lines)

        if suffix == ".py":
            lines = path.read_text("utf-8", errors="ignore").splitlines()
            return "\n".join(
                l.lstrip("# ").strip()
                for l in lines
                if l.strip().startswith("#")
            )

        return ""